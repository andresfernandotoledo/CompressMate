import os
import subprocess
import tempfile
import shutil
import pdfplumber
from flask import Flask, render_template, request, send_file, abort
from io import BytesIO
from PIL import Image, UnidentifiedImageError
import moviepy.editor as mp
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from openpyxl import Workbook, load_workbook
from pdf2docx import Converter as PDF2DOCXConverter
from werkzeug.utils import secure_filename
import magic
from zipfile import ZipFile  
from fpdf import FPDF  


app = Flask(__name__)

def detect_file_type(file_stream, filename):
    file_stream.seek(0)
    mime = magic.Magic(mime=True)
    mime_type = mime.from_buffer(file_stream.read(2048))
    file_stream.seek(0)

    # Definir tipos de archivos basados en MIME type
    mime_map = {
        'image': 'imagen',
        'pdf': 'pdf',
        'vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
        'video': 'video'
    }

    for key, value in mime_map.items():
        if key in mime_type:
            return value

    # Si MIME type no se encuentra, intentar con extensión del archivo
    ext = os.path.splitext(filename)[1].lower()
    ext_map = {
        '.docx': 'docx',
        '.pptx': 'pptx',
        '.xlsx': 'xlsx',
        '.pdf': 'pdf',
        '.jpeg': 'imagen',
        '.jpg': 'imagen',
        '.png': 'imagen',
        '.bmp': 'imagen',
        '.gif': 'imagen',
        '.tiff': 'imagen',
        '.mp4': 'video',
        '.mov': 'video',
        '.avi': 'video',
        '.mkv': 'video'
    }

    return ext_map.get(ext, 'desconocido')


# Funciones de compresión
def compress_image(input_stream, quality=85, output_format='jpeg'):
    try:
        output = BytesIO()
        with Image.open(input_stream) as img:
            img = img.convert("RGB")  # Convertir a RGB si la imagen no está en RGB
            img.save(output, format=output_format.upper(), quality=quality)
        output.seek(0)
        return output
    except Exception as e:
        raise ValueError(f"Ocurrió un error al procesar la imagen: {str(e)}")

def compress_pdf(input_stream, target_size_bytes):
    temp_input_path = tempfile.mktemp(suffix='.pdf')
    temp_output_path = tempfile.mktemp(suffix='.pdf')

    with open(temp_input_path, 'wb') as f:
        f.write(input_stream.read())

    # Ajuste iterativo para alcanzar el tamaño objetivo
    def compress_with_quality(quality):
        gs_command = [
            "gs",
            "-sDEVICE=pdfwrite",
            f"-dPDFSETTINGS={quality}",
            "-dNOPAUSE",
            "-dQUIET",
            "-dBATCH",
            f"-sOutputFile={temp_output_path}",
            temp_input_path
        ]
        subprocess.run(gs_command, check=True)

    # Configuraciones de calidad para probar
    quality_settings = ['/screen', '/ebook', '/printer', '/prepress']
    min_target_size_bytes = 500_000 if os.path.getsize(temp_input_path) > 1_000_000 else target_size_bytes

    for quality in quality_settings:
        compress_with_quality(quality)
        output_size_bytes = os.path.getsize(temp_output_path)
        
        if output_size_bytes <= target_size_bytes:
            break

    # Si aún no se ha alcanzado el tamaño objetivo, prueba una compresión más agresiva
    if output_size_bytes > target_size_bytes:
        lower_quality = '/screen'  # El ajuste más bajo
        compress_with_quality(lower_quality)
        output_size_bytes = os.path.getsize(temp_output_path)

    # Ajuste final si es necesario
    if output_size_bytes > target_size_bytes:
        print("No se pudo comprimir a menos del tamaño deseado.")
    
    with open(temp_output_path, 'rb') as f:
        output_data = f.read()

    os.remove(temp_input_path)
    os.remove(temp_output_path)

    output_stream = BytesIO(output_data)
    output_stream.seek(0)
    return output_stream


def compress_docx(input_stream, target_size_bytes):
    temp_input_path = tempfile.mktemp(suffix='.docx')
    temp_output_path = tempfile.mktemp(suffix='.docx')

    with open(temp_input_path, 'wb') as f:
        f.write(input_stream.read())

    doc = Document(temp_input_path)
    
    # Reduce la calidad de las imágenes en el documento
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img = Image.open(BytesIO(rel.target_part.blob))
            img = img.convert("RGB")
            output = BytesIO()
            img.save(output, format='JPEG', quality=30)  # Ajusta la calidad
            rel.target_part._blob = output.getvalue()

    doc.save(temp_output_path)

    # Comprobamos el tamaño del archivo y aplicamos el límite mínimo
    if os.path.getsize(temp_input_path) > 1_000_000:
        min_target_size_bytes = 500_000
    else:
        min_target_size_bytes = target_size_bytes

    output_size_bytes = os.path.getsize(temp_output_path)
    if output_size_bytes > min_target_size_bytes:
        print("No se pudo comprimir a menos del tamaño deseado.")
    
    with open(temp_output_path, 'rb') as f:
        output_data = f.read()

    os.remove(temp_input_path)
    os.remove(temp_output_path)

    output_stream = BytesIO(output_data)
    output_stream.seek(0)
    return output_stream


def compress_pptx(input_stream, target_size_bytes):
    temp_input_path = tempfile.mktemp(suffix='.pptx')
    temp_output_path = tempfile.mktemp(suffix='.pptx')

    with open(temp_input_path, 'wb') as f:
        f.write(input_stream.read())

    ppt = Presentation(temp_input_path)
    
    # Reduce la calidad de las imágenes en el PPT
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                img = Image.open(BytesIO(shape.image.blob))
                img = img.convert("RGB")
                output = BytesIO()
                img.save(output, format='JPEG', quality=30)  # Ajusta la calidad
                shape.image = output.getvalue()

    ppt.save(temp_output_path)

    # Comprobamos el tamaño del archivo y aplicamos el límite mínimo
    if os.path.getsize(temp_input_path) > 1_000_000:
        min_target_size_bytes = 500_000
    else:
        min_target_size_bytes = target_size_bytes

    output_size_bytes = os.path.getsize(temp_output_path)
    if output_size_bytes > min_target_size_bytes:
        print("No se pudo comprimir a menos del tamaño deseado.")
    
    with open(temp_output_path, 'rb') as f:
        output_data = f.read()

    os.remove(temp_input_path)
    os.remove(temp_output_path)

    output_stream = BytesIO(output_data)
    output_stream.seek(0)
    return output_stream


def compress_xlsx(input_stream, target_size_bytes):
    temp_input_path = tempfile.mktemp(suffix='.xlsx')
    temp_output_path = tempfile.mktemp(suffix='.xlsx')

    with open(temp_input_path, 'wb') as f:
        f.write(input_stream.read())

    workbook = load_workbook(temp_input_path)
    workbook.save(temp_output_path)

    # Comprobamos el tamaño del archivo y aplicamos el límite mínimo
    if os.path.getsize(temp_input_path) > 1_000_000:
        min_target_size_bytes = 500_000
    else:
        min_target_size_bytes = target_size_bytes

    output_size_bytes = os.path.getsize(temp_output_path)
    if output_size_bytes > min_target_size_bytes:
        print("No se pudo comprimir a menos del tamaño deseado.")
    
    with open(temp_output_path, 'rb') as f:
        output_data = f.read()

    os.remove(temp_input_path)
    os.remove(temp_output_path)

    output_stream = BytesIO(output_data)
    output_stream.seek(0)
    return output_stream



def compress_video(input_stream, bitrate="500k"):
    try:
        # Crear un archivo temporal para la entrada
        with tempfile.NamedTemporaryFile(delete=False, suffix='.tmp') as temp_input_file:
            temp_input_path = temp_input_file.name
            temp_input_file.write(input_stream.read())
        
        # Obtener la extensión del archivo original
        input_file_name = 'video'
        input_file_ext = '.mp4'  # Default extension
        
        if '.' in input_file_name:
            input_file_ext = os.path.splitext(input_file_name)[1]  # Extract extension

        # Crear un archivo temporal para la salida
        with tempfile.NamedTemporaryFile(delete=False, suffix=input_file_ext) as temp_output_file:
            temp_output_path = temp_output_file.name
        
        # Procesar el video con MoviePy
        video = mp.VideoFileClip(temp_input_path)
        video_resized = video.resize(height=720)
        video_resized.write_videofile(temp_output_path, bitrate=bitrate, codec='libx264', audio_codec='aac')
        
        # Eliminar archivo temporal de entrada
        os.remove(temp_input_path)
        
        return temp_output_path
    except mp.MoviePyError as e:
        raise ValueError(f"Error al procesar el video con MoviePy: {str(e)}")
    except Exception as e:
        raise ValueError(f"Ocurrió un error al procesar el video: {str(e)}")
    
# Función para convertir imágenes a diferentes formatos
def convert_image(input_stream, output_format):
    valid_formats = ['jpeg', 'png', 'bmp', 'gif', 'tiff']

    if output_format.lower() not in valid_formats:
        raise ValueError("Formato de imagen no válido. Los formatos válidos son: " + ", ".join(valid_formats))
    
    output = BytesIO()

    try:
        with Image.open(input_stream) as img:
            if output_format.lower() in ['jpeg', 'bmp', 'gif', 'tiff'] and img.mode not in ['RGB', 'L']:
                img = img.convert('RGB')
            
            img.save(output, format=output_format.upper())
    except UnidentifiedImageError:
        raise ValueError("El archivo no es una imagen válida.")
    except Exception as e:
        raise ValueError(f"Ocurrió un error al convertir la imagen: {str(e)}")

    output.seek(0)
    return output

# Funciones para convertir archivos PDF a diferentes formatos
def convert_pdf(input_stream, output_format):
    if output_format == 'docx':
        return convert_pdf_to_docx(input_stream)
    elif output_format == 'pptx':
        return convert_pdf_to_pptx(input_stream)
    elif output_format == 'xlsx':
        return convert_pdf_to_xlsx(input_stream)
    elif output_format == 'jpeg':
        return convert_pdf_to_image(input_stream, 'JPEG')
    elif output_format == 'png':
        return convert_pdf_to_image(input_stream, 'PNG')
    else:
        raise ValueError("Formato de salida no soportado para PDF.")

def convert_pdf_to_docx(input_stream):
    output = BytesIO()
    with tempfile.NamedTemporaryFile(delete=True, suffix='.pdf') as temp_input, \
         tempfile.NamedTemporaryFile(delete=True, suffix='.docx') as temp_output:
        temp_input.write(input_stream.read())
        temp_input.flush()

        converter = PDF2DOCXConverter(temp_input.name)
        converter.convert(temp_output.name)
        converter.close()

        with open(temp_output.name, 'rb') as f:
            output.write(f.read())

    output.seek(0)
    return output

def convert_pdf_to_pptx(input_stream):
    output = BytesIO()
    with tempfile.NamedTemporaryFile(delete=True, suffix='.pdf') as temp_input, \
         tempfile.NamedTemporaryFile(delete=True, suffix='.pptx') as temp_output:
        temp_input.write(input_stream.read())
        temp_input.flush()

        presentation = Presentation()
        with pdfplumber.open(temp_input.name) as pdf:
            for page in pdf.pages:
                slide_layout = presentation.slide_layouts[5]
                slide = presentation.slides.add_slide(slide_layout)
                textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
                textbox.text = page.extract_text()

        presentation.save(temp_output.name)
        with open(temp_output.name, 'rb') as f:
            output.write(f.read())

    output.seek(0)
    return output

def convert_pdf_to_xlsx(input_stream):
    output = BytesIO()
    with tempfile.NamedTemporaryFile(delete=True, suffix='.pdf') as temp_input, \
         tempfile.NamedTemporaryFile(delete=True, suffix='.xlsx') as temp_output:
        temp_input.write(input_stream.read())
        temp_input.flush()

        workbook = Workbook()
        sheet = workbook.active
        with pdfplumber.open(temp_input.name) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                for i, line in enumerate(text.split('\n'), start=1):
                    sheet[f'A{i}'] = line

        workbook.save(temp_output.name)
        with open(temp_output.name, 'rb') as f:
            output.write(f.read())

    output.seek(0)
    return output

def convert_pdf_to_image(input_stream, format):
    output = BytesIO()
    with tempfile.NamedTemporaryFile(delete=True, suffix='.pdf') as temp_input:
        temp_input.write(input_stream.read())
        temp_input.flush()

        images = convert_from_path(temp_input.name)
        images[0].save(output, format=format)

    output.seek(0)
    return output

# Funciones para convertir archivos DOCX a diferentes formatos
def convert_docx(input_stream, output_format):
    if output_format == 'pdf':
        return convert_docx_to_pdf(input_stream)
    elif output_format == 'pptx':
        return convert_docx_to_pptx(input_stream)
    elif output_format == 'xlsx':
        return convert_docx_to_xlsx(input_stream)
    else:
        raise ValueError("Formato de salida no soportado para DOCX.")

def convert_docx_to_pdf(input_stream):
    try:
        output = BytesIO()
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_input, \
             tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_output:
            temp_input.write(input_stream.read())
            temp_input.flush()
            
            subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', temp_input.name, '--outdir', tempfile.gettempdir()], check=True)
            
            with open(temp_output.name, 'rb') as f:
                output.write(f.read())

            os.remove(temp_input.name)
            os.remove(temp_output.name)
        
        output.seek(0)
        return output
    except Exception as e:
        raise ValueError(f"Error convirtiendo DOCX a PDF: {str(e)}")

def convert_docx_to_pptx(input_stream):
    try:
        output = BytesIO()
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_input, \
             tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_output:
            temp_input.write(input_stream.read())
            temp_input.flush()

            doc = Document(temp_input.name)
            presentation = Presentation()

            for paragraph in doc.paragraphs:
                slide_layout = presentation.slide_layouts[5]
                slide = presentation.slides.add_slide(slide_layout)
                textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
                textbox.text = paragraph.text

            presentation.save(temp_output.name)

            with open(temp_output.name, 'rb') as f:
                output.write(f.read())

            os.remove(temp_input.name)
            os.remove(temp_output.name)
        
        output.seek(0)
        return output
    except Exception as e:
        raise ValueError(f"Error convirtiendo DOCX a PPTX: {str(e)}")

def convert_docx_to_xlsx(input_stream):
    try:
        output = BytesIO()
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_input, \
             tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_output:
            temp_input.write(input_stream.read())
            temp_input.flush()

            doc = Document(temp_input.name)
            workbook = Workbook()
            sheet = workbook.active

            for i, paragraph in enumerate(doc.paragraphs, start=1):
                sheet[f'A{i}'] = paragraph.text

            workbook.save(temp_output.name)

            with open(temp_output.name, 'rb') as f:
                output.write(f.read())

            os.remove(temp_input.name)
            os.remove(temp_output.name)
        
        output.seek(0)
        return output
    except Exception as e:
        raise ValueError(f"Error convirtiendo DOCX a XLSX: {str(e)}")

# Funciones para convertir archivos PPTX a diferentes formatos
def convert_pptx(input_stream, output_format):
    if output_format == 'pdf':
        return convert_pptx_to_pdf(input_stream)
    elif output_format == 'docx':
        return convert_pptx_to_docx(input_stream)
    elif output_format == 'xlsx':
        return convert_pptx_to_xlsx(input_stream)
    else:
        raise ValueError("Formato de salida no soportado para PPTX.")

def convert_pptx_to_pdf(input_stream):
    try:
        output = BytesIO()
        pdf = FPDF()
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_input, \
             tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_output:
            temp_input.write(input_stream.read())
            temp_input.flush()

            presentation = Presentation(temp_input.name)
            
            for slide in presentation.slides:
                pdf.add_page()
                pdf.set_font("Arial", size=12)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        pdf.multi_cell(0, 10, shape.text)
            
            pdf.output(temp_output.name)

            with open(temp_output.name, 'rb') as f:
                output.write(f.read())

            os.remove(temp_input.name)
            os.remove(temp_output.name)
        
        output.seek(0)
        return output
    except Exception as e:
        raise ValueError(f"Error convirtiendo PPTX a PDF: {str(e)}")

def convert_pptx_to_docx(input_stream):
    output = BytesIO()
    with tempfile.NamedTemporaryFile(delete=True, suffix='.pptx') as temp_input, \
         tempfile.NamedTemporaryFile(delete=True, suffix='.docx') as temp_output:
        temp_input.write(input_stream.read())
        temp_input.flush()

        presentation = Presentation(temp_input.name)
        doc = Document()

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    doc.add_paragraph(shape.text)

        doc.save(temp_output.name)
        with open(temp_output.name, 'rb') as f:
            output.write(f.read())

    output.seek(0)
    return output

def convert_pptx_to_xlsx(input_stream):
    output = BytesIO()
    with tempfile.NamedTemporaryFile(delete=True, suffix='.pptx') as temp_input, \
         tempfile.NamedTemporaryFile(delete=True, suffix='.xlsx') as temp_output:
        temp_input.write(input_stream.read())
        temp_input.flush()

        presentation = Presentation(temp_input.name)
        workbook = Workbook()
        sheet = workbook.active

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    sheet.append([shape.text])

        workbook.save(temp_output.name)
        with open(temp_output.name, 'rb') as f:
            output.write(f.read())

    output.seek(0)
    return output

# Funciones para convertir archivos XLSX a diferentes formatos
def convert_xlsx(input_stream, output_format):
    if output_format == 'pdf':
        return convert_xlsx_to_pdf(input_stream)
    elif output_format == 'pptx':
        return convert_xlsx_to_pptx(input_stream)
    elif output_format == 'docx':
        return convert_xlsx_to_docx(input_stream)
    else:
        raise ValueError("Formato de salida no soportado para XLSX.")

def convert_xlsx_to_pdf(input_stream):
    try:
        output = BytesIO()
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_input, \
             tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_output:
            temp_input.write(input_stream.read())
            temp_input.flush()

            workbook = load_workbook(temp_input.name)
            pdf = FPDF()

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                pdf.add_page()
                pdf.set_font("Arial", size=12)
                pdf.cell(200, 10, txt=sheet_name, ln=True, align='C')

                for row in sheet.iter_rows(values_only=True):
                    line = "\t".join(map(str, row))
                    pdf.multi_cell(0, 10, line)

            pdf.output(temp_output.name)

            with open(temp_output.name, 'rb') as f:
                output.write(f.read())

            os.remove(temp_input.name)
            os.remove(temp_output.name)
        
        output.seek(0)
        return output
    except Exception as e:
        raise ValueError(f"Error convirtiendo XLSX a PDF: {str(e)}")

def convert_xlsx_to_pptx(input_stream):
    output = BytesIO()
    with tempfile.NamedTemporaryFile(delete=True, suffix='.xlsx') as temp_input, \
         tempfile.NamedTemporaryFile(delete=True, suffix='.pptx') as temp_output:
        
        # Guardar el contenido del archivo XLSX en un archivo temporal
        temp_input.write(input_stream.read())
        temp_input.flush()

        # Cargar el archivo XLSX
        wb = load_workbook(temp_input.name)

        # Crear una presentación PPTX
        presentation = Presentation()

        # Iterar sobre cada hoja en el archivo XLSX
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            # Crear una diapositiva para cada hoja
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
            text_frame = textbox.text_frame

            # Ajustar el tamaño y estilo del texto
            p = text_frame.add_paragraph()
            p.text = f"Hoja: {sheet_name}"
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 0, 0)  # Cambiar color si es necesario

            # Agregar datos de la hoja de cálculo a la diapositiva
            for row in sheet.iter_rows(values_only=True):
                # Crear un párrafo para cada fila de datos
                p = text_frame.add_paragraph()
                # Combinar las celdas de la fila en una sola cadena de texto, omitiendo los valores None
                p.text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                p.font.size = Pt(12)  # Ajustar el tamaño de la fuente del texto de datos

        # Guardar la presentación en el archivo temporal
        presentation.save(temp_output.name)

        # Leer el archivo PPTX y devolverlo como un BytesIO stream
        with open(temp_output.name, 'rb') as f:
            output.write(f.read())

    output.seek(0)
    return output



def convert_xlsx_to_docx(input_stream):
    output = BytesIO()
    with tempfile.NamedTemporaryFile(delete=True, suffix='.xlsx') as temp_input, \
         tempfile.NamedTemporaryFile(delete=True, suffix='.docx') as temp_output:
        
        # Guardar el contenido del archivo XLSX en un archivo temporal
        temp_input.write(input_stream.read())
        temp_input.flush()

        # Cargar el archivo XLSX
        wb = load_workbook(temp_input.name)

        # Crear un documento DOCX
        doc = Document()

        # Iterar sobre cada hoja en el archivo XLSX
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]

            # Agregar un encabezado para cada hoja
            doc.add_heading(f'Hoja: {sheet_name}', level=1)

            # Agregar datos de la hoja de cálculo al documento
            for row in sheet.iter_rows(values_only=True):
                # Crear un párrafo para cada fila de datos
                p = doc.add_paragraph()
                # Combinar las celdas de la fila en una sola cadena de texto, omitiendo los valores None
                run = p.add_run("\t".join([str(cell) if cell is not None else "" for cell in row]))
                run.font.size = Pt(12)

        # Guardar el documento en el archivo temporal
        doc.save(temp_output.name)

        # Leer el archivo DOCX y devolverlo como un BytesIO stream
        with open(temp_output.name, 'rb') as f:
            output.write(f.read())

    output.seek(0)
    return output

# Función principal de conversión
def convert(input_stream, input_type, output_format):
    if input_type == 'pdf':
        return convert_pdf(input_stream, output_format)
    elif input_type == 'docx':
        return convert_docx(input_stream, output_format)
    elif input_type == 'pptx':
        return convert_pptx(input_stream, output_format)
    elif input_type == 'xlsx':
        return convert_xlsx(input_stream, output_format)
    else:
        raise ValueError("Tipo de entrada no soportado")

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        file = request.files.get('file')
        if not file:
            abort(400, description="No se ha enviado ningún archivo.")
        
        action_type = request.form.get('action_type')
        output_format = request.form.get('output_format', '').lower()
        quality = request.form.get('quality', 85, type=int)
        quality_kb = request.form.get('quality_kb', 100, type=int)
        bitrate = request.form.get('bitrate', '500k')
        file_stream = BytesIO(file.read())
        filename = secure_filename(file.filename)

        detected_file_type = detect_file_type(file_stream, filename)
        file_stream.seek(0)

        if detected_file_type == 'desconocido':
            abort(400, description="Tipo de archivo no soportado o no reconocido.")

        try:
            if action_type == 'compress':
                return handle_compression(file_stream, detected_file_type, output_format, quality, quality_kb, bitrate)
            elif action_type == 'convert':
                return handle_conversion(file_stream, detected_file_type, output_format, filename)
            else:
                raise ValueError("Acción no soportada.")
        except ValueError as e:
            abort(400, description=str(e))
        except Exception as e:
            import traceback
            traceback.print_exc()
            abort(500, description=f"Error interno del servidor: {str(e)}")

    return render_template('index.html')

def get_image_format(file_stream):
    try:
        with Image.open(file_stream) as img:
            return img.format.lower()
    except Exception as e:
        raise ValueError(f"No se pudo detectar el formato de la imagen: {str(e)}")

def handle_compression(file_stream, file_type, output_format, quality, quality_kb, bitrate):
    try:
        if file_type == 'imagen':
            # Detecta el formato original
            file_stream.seek(0)  # Asegúrate de que el stream esté en el inicio
            original_format = get_image_format(file_stream)
            
            # Si el formato original no está en la lista, usa 'jpeg' por defecto
            valid_formats = ['jpeg', 'png', 'bmp', 'gif', 'tiff']
            if original_format not in valid_formats:
                original_format = 'jpeg'

            # Reposiciona el stream para la compresión
            file_stream.seek(0)
            
            output_file = compress_image(file_stream, quality, original_format)
            
            mime_type = {
                'jpeg': 'image/jpeg',
                'png': 'image/png',
                'bmp': 'image/bmp',
                'gif': 'image/gif',
                'tiff': 'image/tiff'
            }.get(original_format, 'application/octet-stream')
            
            filename = f'compressed_file.{original_format}'
            return send_file(output_file, as_attachment=True, download_name=filename, mimetype=mime_type)

        
        
        elif file_type == 'pdf':
            output_file = compress_pdf(file_stream, quality_kb)
            return send_file(output_file, as_attachment=True, download_name='compressed_file.pdf', mimetype='application/pdf')

        elif file_type == 'docx':
            output_file = compress_docx(file_stream, quality_kb)  # Supongamos que tienes una función compress_docx
            return send_file(output_file, as_attachment=True, download_name='compressed_file.docx', mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')

        elif file_type == 'pptx':
            output_file = compress_pptx(file_stream, quality_kb)  # Supongamos que tienes una función compress_pptx
            return send_file(output_file, as_attachment=True, download_name='compressed_file.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        elif file_type == 'xlsx':
            output_file = compress_xlsx(file_stream, quality_kb)  # Supongamos que tienes una función compress_xlsx
            return send_file(output_file, as_attachment=True, download_name='compressed_file.xlsx', mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        
        elif file_type == 'video':
            # Suponiendo que la extensión original se pasa como parte del input
            original_filename = 'video.mp4'  # Esto debería ser ajustado según cómo obtienes el nombre del archivo original.
            input_ext = os.path.splitext(original_filename)[1]  # Obtener la extensión del archivo original
            
            output_file_path = compress_video(file_stream, bitrate)
            
            # Usar la extensión original para el nombre del archivo de salida
            filename = f'compressed_video{input_ext}'
            return send_file(output_file_path, as_attachment=True, download_name=filename, mimetype=f'video{input_ext}')

        else:
            raise ValueError("Tipo de archivo no soportado para compresión.")
    except ValueError as e:
        abort(400, description=str(e))
    except Exception as e:
        abort(500, description=f"Error interno del servidor: {str(e)}")

def handle_conversion(file_stream, file_type, output_format, filename=None):
    output_format = output_format.lower()

    if not filename:
        abort(400, description="Nombre de archivo no proporcionado.")
    
    filename = secure_filename(filename)

    valid_formats = {
        'imagen': ['jpeg', 'png', 'bmp', 'gif', 'tiff'],
        'pdf': ['docx', 'pptx', 'xlsx', 'jpeg', 'png'],
        'docx': ['pdf', 'pptx', 'xlsx'],
        'pptx': ['pdf', 'docx', 'xlsx'],
        'xlsx': ['pdf', 'docx', 'pptx'],
    }

    if file_type not in valid_formats:
        abort(400, description=f"Tipo de archivo no soportado: {file_type}.")
    
    if output_format not in valid_formats.get(file_type, []):
        abort(400, description=f"Formato de salida no válido para {file_type}. Los formatos válidos son: {', '.join(valid_formats[file_type])}.")

    try:
        if file_type == 'imagen':
            output_file = convert_image(file_stream, output_format)
            mime_type = f'image/{output_format}'
        else:
            output_file = convert(file_stream, file_type, output_format)
            mime_type = {
                'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                'pdf': 'application/pdf',
                'jpeg': 'image/jpeg',
                'png': 'image/png'
            }.get(output_format, 'application/octet-stream')

        output_filename = f'converted_file.{output_format}'
        return send_file(output_file, as_attachment=True, download_name=output_filename, mimetype=mime_type)
    except ValueError as e:
        abort(400, description=str(e))
    except Exception as e:
        import traceback
        traceback.print_exc()
        abort(500, description=f"Error interno del servidor: {str(e)}")
        
if __name__ == '__main__':
    app.run(debug=True)
